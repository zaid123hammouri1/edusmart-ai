# routes/file_routes.py - نظام رفع الملفات الحقيقي
"""
Real File Upload System:
1. Actual file upload to server
2. Text extraction from PDF/DOCX
3. Content stored for RAG chatbot
"""
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.orm import Session
from typing import Optional
import os
import shutil
from datetime import datetime
from pathlib import Path
from database import get_db
from auth import get_lecturer, get_current_user
import models

router = APIRouter(prefix="/files", tags=["Files"])

# Upload directory
UPLOAD_DIR = Path(__file__).parent.parent / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)

# Allowed file types
ALLOWED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".pptx", ".xlsx"}
MAX_FILE_SIZE = 10 * 1024 * 1024  # 10MB


def extract_text_from_file(file_path: Path) -> str:
    """Extract text content from uploaded file for RAG"""
    text = ""
    suffix = file_path.suffix.lower()
    
    try:
        if suffix == ".txt":
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        
        elif suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(str(file_path))
                for page in doc:
                    text += page.get_text()
                doc.close()
            except ImportError:
                # Fallback: try pdfplumber
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text += page_text + "\n"
                except ImportError:
                    text = "[PDF text extraction not available - install PyMuPDF or pdfplumber]"
        
        elif suffix in {".docx", ".doc"}:
            try:
                import docx
                doc = docx.Document(str(file_path))
                text = "\n".join([para.text for para in doc.paragraphs])
            except ImportError:
                text = "[DOCX text extraction not available - install python-docx]"
        
        elif suffix == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            except ImportError:
                text = "[PPTX text extraction not available - install python-pptx]"
    
    except Exception as e:
        text = f"[Error extracting text: {str(e)}]"
    
    return text.strip()


@router.post("/upload/{course_id}")
async def upload_course_material(
    course_id: int,
    file: UploadFile = File(...),
    title: str = Form(...),
    description: str = Form(""),
    current_user: models.User = Depends(get_lecturer),
    db: Session = Depends(get_db)
):
    """Upload a file as course material (Lecturer only)"""
    
    # Verify lecturer owns this course
    course = db.query(models.Course).filter(
        models.Course.id == course_id,
        models.Course.lecturer_id == current_user.id
    ).first()
    
    if not course:
        raise HTTPException(status_code=403, detail="You don't teach this course")
    
    # Validate file
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file provided")
    
    file_ext = Path(file.filename).suffix.lower()
    if file_ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400, 
            detail=f"File type not allowed. Allowed: {', '.join(ALLOWED_EXTENSIONS)}"
        )
    
    # Check file size
    file.file.seek(0, 2)  # Seek to end
    file_size = file.file.tell()
    file.file.seek(0)  # Reset
    
    if file_size > MAX_FILE_SIZE:
        raise HTTPException(status_code=400, detail="File too large (max 10MB)")
    
    # Create course folder
    course_folder = UPLOAD_DIR / f"course_{course_id}"
    course_folder.mkdir(exist_ok=True)
    
    # Generate unique filename
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_filename = f"{timestamp}_{file.filename}"
    file_path = course_folder / safe_filename
    
    # Save file
    try:
        with open(file_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
    
    # Extract text for RAG
    extracted_text = extract_text_from_file(file_path)
    
    # Save to database
    material = models.CourseMaterial(
        course_id=course_id,
        title=title,
        description=description,
        file_url=str(file_path),
        content_text=extracted_text[:10000] if extracted_text else None,  # Limit stored text
        uploaded_by=current_user.id
    )
    db.add(material)
    db.commit()
    db.refresh(material)
    
    return {
        "message": "File uploaded successfully",
        "material_id": material.id,
        "file_name": file.filename,
        "text_extracted": len(extracted_text) > 0,
        "text_length": len(extracted_text)
    }


@router.get("/course/{course_id}")
def get_course_files(
    course_id: int,
    current_user: models.User = Depends(get_current_user),
    db: Session = Depends(get_db)
):
    """Get all files for a course"""
    
    # Check access (student enrolled or lecturer owns)
    if current_user.role == "student":
        enrollment = db.query(models.Enrollment).filter(
            models.Enrollment.student_id == current_user.id,
            models.Enrollment.course_id == course_id
        ).first()
        if not enrollment:
            raise HTTPException(status_code=403, detail="Not enrolled in this course")
    
    elif current_user.role == "lecturer":
        course = db.query(models.Course).filter(
            models.Course.id == course_id,
            models.Course.lecturer_id == current_user.id
        ).first()
        if not course:
            raise HTTPException(status_code=403, detail="You don't teach this course")
    
    # Get materials
    materials = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.course_id == course_id
    ).all()
    
    return [
        {
            "id": m.id,
            "title": m.title,
            "description": m.description,
            "file_name": Path(m.file_url).name if m.file_url else None,
            "has_text": bool(m.content_text),
            "created_at": str(m.created_at) if m.created_at else None
        }
        for m in materials
    ]


@router.delete("/{material_id}")
def delete_material(
    material_id: int,
    current_user: models.User = Depends(get_lecturer),
    db: Session = Depends(get_db)
):
    """Delete a material (Lecturer only)"""
    
    material = db.query(models.CourseMaterial).filter(
        models.CourseMaterial.id == material_id
    ).first()
    
    if not material:
        raise HTTPException(status_code=404, detail="Material not found")
    
    # Verify ownership
    course = db.query(models.Course).filter(
        models.Course.id == material.course_id,
        models.Course.lecturer_id == current_user.id
    ).first()
    
    if not course:
        raise HTTPException(status_code=403, detail="You don't own this material")
    
    # Delete file from disk
    if material.file_url and Path(material.file_url).exists():
        try:
            Path(material.file_url).unlink()
        except:
            pass
    
    # Delete from database
    db.delete(material)
    db.commit()
    
    return {"message": "Material deleted"}
